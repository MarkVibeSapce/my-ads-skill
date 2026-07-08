#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
gen_report.py — ออกรายงานผลแอดส่ง "หัวหน้างาน" (weekly / monthly)
เป็น 3 ไฟล์: HTML + PDF + PPTX · ดีไซน์เดียวกัน · มาตรฐานเดียวกับรายงานเอเจนซี่
(header ribbon เข้ม + accent CI ของแบรนด์ + KPI card + cover + สรุปผู้บริหาร)

แนวคิด: แยก "ตัวเลข" (ดึงจาก API) ออกจาก "บทวิเคราะห์" (Claude เขียน)
  1) FETCH  → ดึงตัวเลขจริงจาก Meta API เขียนลง report_data.json (ช่องวิเคราะห์ว่างไว้)
  2) Claude เติมช่องวิเคราะห์ (exec/goods/improves/next) ใน report_data.json
  3) RENDER → อ่าน json ออกเป็น HTML + PDF + PPTX

รันง่ายสุด (ดูตัวอย่างทันที ไม่ต้องต่อ API):
  python3 gen_report.py --period monthly --mock --out all
  python3 gen_report.py --period weekly  --mock --out all

ใช้จริง:
  python3 gen_report.py --period weekly  --fetch                 # ดึงเลข → report_data.json
  (Claude เติมบทวิเคราะห์ใน report_data.json)
  python3 gen_report.py --period weekly  --render --out all      # ออก 3 ไฟล์

ตั้ง CI แบรนด์ (ใน .env — รายงานจะออกสีของแบรนด์คุณเอง):
  BRAND_NAME=บ้านในฝัน ดีไซน์
  BRAND_TAGLINE=รับออกแบบและสร้างบ้าน
  BRAND_ACCENT=2E5BE6      # สีหลักแบรนด์ (hex ไม่มี #) — ใช้เน้นค่า/เส้น/ปุ่ม
  BRAND_DARK=16202E        # สี ribbon หัวสไลด์ (เข้ม)
  CPR_TARGET=150           # เป้า CPR — ใช้ตัดสถานะ ✅/🔴
  ACCOUNT_ID=act_xxx       # จาก BRAND_PROFILE
  REPORT_EN_FONT=Sarabun   # ฟอนต์ตัวเลข (default Sarabun · ตั้ง Montserrat ได้ถ้าเครื่องมี)
"""
import sys, os, json, argparse, subprocess, urllib.request, urllib.parse

try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

HERE = os.path.dirname(os.path.abspath(__file__))
SKILL_DIR = os.path.dirname(HERE)
DATA_JSON = os.path.join(HERE, "report_data.json")
API = "https://graph.facebook.com/v21.0"

# ---------- .env ----------
def load_env():
    env = {}
    p = os.path.join(SKILL_DIR, ".env")
    if os.path.exists(p):
        for line in open(p, encoding="utf-8"):
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                k, v = line.split("=", 1)
                env[k.strip()] = v.strip()
    return env

def load_brand_profile():
    """อ่านค่า KEY=value / KEY = value จาก BRAND_PROFILE.md (ACCOUNT_ID/PAGE_ID/BRAND_NAME...)
    ให้นักเรียนกรอกไฟล์เดียว (BRAND_PROFILE) แล้วรายงานใช้ได้ — ไม่ต้องกรอก .env ซ้ำ"""
    bp = {}
    p = os.path.join(SKILL_DIR, "BRAND_PROFILE.md")
    if not os.path.exists(p):
        return bp
    keys = {"ACCOUNT_ID", "PAGE_ID", "BRAND_NAME", "BRAND_TAGLINE",
            "BRAND_ACCENT", "BRAND_DARK", "CPR_TARGET", "REPORT_EN_FONT"}
    for line in open(p, encoding="utf-8"):
        line = line.strip().lstrip("|").strip()
        if "=" not in line:
            continue
        k, v = line.split("=", 1)
        k = k.strip()
        v = v.split("#")[0].split("←")[0].strip().strip("`").strip()
        # ข้าม placeholder ที่ยังไม่แก้ (act_XXXX / XXXXX)
        if k in keys and v and "XXX" not in v.upper():
            bp.setdefault(k, v)
    return bp

ENV = load_env()
BP = load_brand_profile()
def cfg(k, default=None):
    # ลำดับ: .env > environment > BRAND_PROFILE.md > default
    return ENV.get(k) or os.environ.get(k) or BP.get(k) or default

# ---------- brand CI (ตั้งสีเองได้ต่อแบรนด์) ----------
def _hx(s):
    s = s.lstrip("#")
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
def _hs(t):
    return "%02X%02X%02X" % t
def _mix(a, b, t):  # a,b hex → mix (t=0 → a, t=1 → b)
    ca, cb = _hx(a), _hx(b)
    return _hs(tuple(round(ca[i] + (cb[i] - ca[i]) * t) for i in range(3)))
def _dark(a, t=0.22):  return _mix(a, "000000", t)
def _soft(a, t=0.90):  return _mix(a, "FFFFFF", t)

ACCENT     = (cfg("BRAND_ACCENT", "2E5BE6") or "2E5BE6").lstrip("#")
ACCENT_INK = _dark(ACCENT, 0.22)
ACCENT_SOFT = _soft(ACCENT, 0.90)
DARK       = (cfg("BRAND_DARK", "16202E") or "16202E").lstrip("#")
DARK_SOFT  = _mix(DARK, "FFFFFF", 0.14)
INK, MUTED, FAINT = "16202E", "5D6B7E", "8A97A8"
LINE, LINE_SOFT = "E4E8EE", "EEF1F5"
PAPER, SURFACE, SURFACE2 = "F6F7F9", "FFFFFF", "FBFCFD"
GOOD, GOOD_SOFT = "16915A", "E4F3EC"
WARN, WARN_SOFT = "B7791F", "F6EDDD"
BAD, BAD_SOFT = "D14343", "FBEAEA"
EN_FONT = cfg("REPORT_EN_FONT", "Sarabun")
TH_FONT = "Sarabun"

def api_get(path, params, token):
    params = dict(params); params["access_token"] = token
    url = f"{API}/{path}?{urllib.parse.urlencode(params)}"
    with urllib.request.urlopen(url) as r:
        return json.loads(r.read())

# ---------- helpers ----------
MSG_STARTED = "onsite_conversion.messaging_conversation_started_7d"
def _find(items, key):
    for x in (items or []):
        if x.get("action_type") == key:
            return float(x.get("value", 0))
    return None

def baht(n):
    return "—" if n is None else f"{int(round(n)):,}"

def pct_change(now, prev):
    if not prev or now is None:
        return None
    return round((now - prev) / prev * 100, 1)

# ================================================================
# FETCH (ตัวเลขจริง)
# ================================================================
def fetch(period, since, until):
    token = cfg("FB_ACCESS_TOKEN_ADS")
    acct = cfg("ACCOUNT_ID")
    if not token or not acct:
        sys.exit("❌ ไม่เจอ FB_ACCESS_TOKEN_ADS หรือ ACCOUNT_ID — เติม .env ก่อน (ดู SETUP.md) "
                 "หรือใช้ --mock เพื่อดูตัวอย่าง")

    def insights(s, u, level="account"):
        return api_get(f"{acct}/insights", {
            "fields": "spend,reach,impressions,cpm,frequency,actions,cost_per_action_type",
            "time_range": json.dumps({"since": s, "until": u}),
            "level": level, "limit": 200,
        }, token).get("data", [])

    now = insights(since, until)
    row = now[0] if now else {}
    spend = float(row.get("spend", 0)) if row else 0
    msgs = _find(row.get("actions"), MSG_STARTED) or 0
    cpr = _find(row.get("cost_per_action_type"), MSG_STARTED)

    prev = _prev_range(since, until)
    prow = (insights(*prev) or [{}])[0]
    msgs_prev = _find(prow.get("actions"), MSG_STARTED) or 0
    cpr_prev = _find(prow.get("cost_per_action_type"), MSG_STARTED)

    personas = []
    for a in insights(since, until, level="adset"):
        pc = _find(a.get("cost_per_action_type"), MSG_STARTED)
        pm = _find(a.get("actions"), MSG_STARTED) or 0
        personas.append({"name": a.get("adset_name", "—"), "meta": "",
                         "spend": float(a.get("spend", 0)), "msgs": int(pm),
                         "cpr": pc, "status": _status(pc)})
    personas.sort(key=lambda x: (x["cpr"] is None, x["cpr"] or 9e9))

    target = float(cfg("CPR_TARGET", 150))
    return {
        "period": period, "brand": cfg("BRAND_NAME", "ธุรกิจของคุณ"),
        "tagline": cfg("BRAND_TAGLINE", "รายงานผลโฆษณา Facebook"),
        "period_label": _period_label(period, since, until),
        "range_label": f"{since} ถึง {until}", "cpr_target": target,
        "spend": spend, "budget": None,
        "msgs": int(msgs), "msgs_prev": int(msgs_prev),
        "cpr": cpr, "cpr_prev": cpr_prev,
        "reach": int(float(row.get("reach", 0))) if row else 0,
        "cpm": float(row.get("cpm", 0)) if row else 0,
        "frequency": round(float(row.get("frequency", 0)), 1) if row else 0,
        "leads": None, "lead_rate": None, "personas": personas,
        "exec_summary": "", "funnel": [], "questions": [],
        "goods": [], "improves": [], "next_plan": "",
        "appointments": None, "closes": None,
        "_note": "เติม exec_summary / funnel / questions / goods / improves / next_plan ก่อน --render",
    }

def _prev_range(since, until):
    from datetime import date, timedelta
    a = date(*map(int, since.split("-"))); b = date(*map(int, until.split("-")))
    span = (b - a).days + 1
    return (a - timedelta(days=span)).isoformat(), (a - timedelta(days=1)).isoformat()

def _status(cpr):
    if cpr is None: return "no-data"
    t = float(cfg("CPR_TARGET", 150))
    if cpr <= t: return "good"
    if cpr <= t * 1.5: return "warn"
    return "bad"

def _period_label(period, since, until):
    th = ["", "ม.ค.", "ก.พ.", "มี.ค.", "เม.ย.", "พ.ค.", "มิ.ย.",
          "ก.ค.", "ส.ค.", "ก.ย.", "ต.ค.", "พ.ย.", "ธ.ค."]
    if period == "monthly":
        y, m, _ = since.split("-"); return f"ประจำเดือน {th[int(m)]} {y}"
    s, u = since.split("-"), until.split("-")
    return f"สัปดาห์ {int(s[2])} {th[int(s[1])]} – {int(u[2])} {th[int(u[1])]} {u[0]}"

# ================================================================
# MOCK (ตัวอย่าง "บ้านในฝัน")
# ================================================================
def mock(period):
    base = {
        "brand": "บ้านในฝัน ดีไซน์", "tagline": "รับออกแบบและสร้างบ้าน", "cpr_target": 150,
        "personas": [
            {"name": "คุณนุ่น — มีที่ดิน", "meta": "35–50 ปี · กทม.+ปริมณฑล",
             "spend": 6240, "msgs": 58, "cpr": 108, "status": "good"},
            {"name": "คุณต้น — คู่แต่งใหม่", "meta": "28–38 ปี · กทม.+ปริมณฑล",
             "spend": 5180, "msgs": 39, "cpr": 133, "status": "good"},
            {"name": "ทดสอบ — กำลังหาที่ดิน", "meta": "30–45 ปี · มุมใหม่",
             "spend": 3400, "msgs": 15, "cpr": 227, "status": "bad"},
        ], "_mock": True,
    }
    if period == "monthly":
        base.update({
            "period": "monthly", "period_label": "ประจำเดือน มิ.ย. 2026",
            "range_label": "1–30 มิ.ย. 2026",
            "spend": 14820, "budget": 15000, "msgs": 112, "msgs_prev": 100,
            "cpr": 132, "cpr_prev": 151, "reach": 84200, "cpm": 62, "frequency": 2.4,
            "leads": 68, "lead_rate": 61,
            "exec_summary": "เดือนนี้แอดทำงาน**ดีกว่ามาตรฐาน** — ได้แชท 112 คน เป็นเลดจริง 68 คน "
                            "(61% สูงกว่าค่าเฉลี่ยธุรกิจรับสร้างบ้าน ~45%) · ค่าต่อแชท ฿132 ต่ำกว่าเป้า "
                            "และดีขึ้นต่อเนื่อง · จุดที่ต้องแก้คือกลุ่มทดสอบที่ยังแพงและคนทักแล้วเงียบบางส่วน",
            "funnel": [["แชททั้งหมด", 112], ["เลดจริง (คุยโต้ตอบ)", 68],
                       ["ถามราคา/ขอแบบ", 44], ["ขอนัด/ดูแบบ", 18]],
            "questions": [["ราคา / งบ", 34], ["แบบบ้าน", 26], ["ที่ดิน", 18],
                          ["สินเชื่อ", 14], ["ระยะเวลา", 8]],
            "goods": ["**คุณภาพแชทสูง** — เลดจริง 61% เหนือค่าเฉลี่ยธุรกิจ",
                      "**CPR ฿132 ต่ำกว่าเป้า** ดีขึ้นต่อเนื่อง 4 สัปดาห์",
                      "กลุ่ม **\"คุณนุ่น\" คุ้มสุด ฿108** — ช่องทางหลักชัด"],
            "improves": ["**กลุ่มทดสอบ ฿227** แพงเกินเป้า กินงบ 23%",
                         "**ทักแล้วเงียบ 44 คน** — ควรตอบใน 5 นาที + ข้อความต้อนรับ",
                         "**34% ถามราคา** แต่โพสต์ไม่มีราคา — ใส่ช่วงราคาเริ่มต้น"],
            "next_plan": "ย้ายงบจากกลุ่มทดสอบ → \"คุณนุ่น\" · ตั้งข้อความตอบอัตโนมัติบอกราคาเริ่มต้น · "
                         "ทำโพสต์ \"ไม่มีที่ดินก็เริ่มได้\" จับกลุ่มที่ถามเรื่องที่ดิน 18%",
            "appointments": 18, "closes": None,
        })
    else:
        base.update({
            "period": "weekly", "period_label": "สัปดาห์ 29 มิ.ย. – 5 ก.ค. 2026",
            "range_label": "29 มิ.ย.–5 ก.ค. 2026",
            "spend": 3480, "budget": 3500, "msgs": 27, "msgs_prev": 22,
            "cpr": 129, "cpr_prev": 151, "reach": 21400, "cpm": 60, "frequency": 1.8,
            "leads": 17, "lead_rate": 63,
            "exec_summary": "สัปดาห์นี้ดีขึ้น — CPR ลดจาก ฿151 → ฿129 (ถูกลง 15%) หลังเปลี่ยนรูปใหม่ · "
                            "แชทเพิ่มเป็น 27 คน · กลุ่ม \"คุณนุ่น\" ยังคุ้มสุด",
            "funnel": [], "questions": [],
            "goods": ["**CPR ลด 15%** เพราะเปลี่ยนรูป \"บ้านเสร็จจริง\"",
                      "**แชท ▲ 23%** จากสัปดาห์ก่อน (22 → 27)"],
            "improves": ["**Frequency 1.8x** เริ่มขยับ — เตรียมรูปสำรอง",
                         "กลุ่มทดสอบยังแพง — เฝ้าดูอีก 3 วัน"],
            "next_plan": "เพิ่มงบกลุ่ม \"คุณนุ่น\" · เตรียมรูปชุดใหม่กันคนเบื่อ",
            "appointments": None, "closes": None,
        })
    return base

# ================================================================
# HTML
# ================================================================
def _css():
    return f"""
:root{{--ink:#{INK};--muted:#{MUTED};--faint:#{FAINT};--line:#{LINE};--line-soft:#{LINE_SOFT};
--paper:#{PAPER};--surface:#{SURFACE};--surface-2:#{SURFACE2};--accent:#{ACCENT};--accent-ink:#{ACCENT_INK};
--accent-soft:#{ACCENT_SOFT};--dark:#{DARK};--dark-soft:#{DARK_SOFT};--good:#{GOOD};--good-soft:#{GOOD_SOFT};
--warn:#{WARN};--warn-soft:#{WARN_SOFT};--bad:#{BAD};--bad-soft:#{BAD_SOFT};
--sans:"{TH_FONT}","TH Sarabun New","Sukhumvit Set",-apple-system,"Segoe UI",system-ui,sans-serif;
color-scheme:light;}}
*{{box-sizing:border-box}}
html,body{{margin:0;background:var(--paper);color:var(--ink);font-family:var(--sans);font-size:14px;
line-height:1.5;-webkit-font-smoothing:antialiased;-webkit-print-color-adjust:exact;print-color-adjust:exact}}
.num{{font-variant-numeric:tabular-nums}}
.label{{font-size:10px;font-weight:700;letter-spacing:.08em;text-transform:uppercase;color:var(--faint)}}
.bar{{position:sticky;top:0;z-index:5;display:flex;align-items:center;gap:12px;padding:10px 16px;
background:var(--surface);border-bottom:1px solid var(--line)}}
.bar p{{margin:0;font-size:12.5px;color:var(--muted)}} .bar .sp{{margin-left:auto}}
.pbtn{{font-family:inherit;font-size:13px;font-weight:700;padding:9px 16px;border:none;border-radius:9px;
background:var(--accent);color:#fff;cursor:pointer}}
.page{{width:210mm;min-height:297mm;margin:16px auto;padding:0 0 8mm;background:var(--surface);
box-shadow:0 2px 18px rgba(20,32,46,.10);overflow:hidden}}
/* ribbon header (มาตรฐานเอเจนซี่: แถบเข้ม + accent) */
.ribbon{{background:var(--dark);color:#fff;padding:14px 13mm 13px;display:flex;align-items:center;gap:13px;
border-bottom:3px solid var(--accent)}}
.rmark{{width:46px;height:46px;border-radius:11px;flex-shrink:0;
background:linear-gradient(145deg,var(--accent),var(--accent-ink));
display:flex;align-items:center;justify-content:center;color:#fff;font-weight:800;font-size:21px}}
.ribbon h1{{margin:0;font-size:21px;font-weight:750;letter-spacing:-.01em;color:#fff}}
.ribbon .rt{{margin:2px 0 0;font-size:12.5px;color:#C7D0DD}}
.ribbon .rmeta{{margin-left:auto;text-align:right;font-size:11px;color:#AAB4C4;line-height:1.7}}
.ribbon .rmeta b{{color:#E6EBF2;font-weight:650}}
.ribbon .rmeta .chip{{display:inline-block;margin-top:3px;padding:2px 9px;border-radius:20px;
background:var(--accent);color:#fff;font-weight:700;font-size:10.5px}}
.body{{padding:0 13mm}}
h2.sec{{font-size:14px;font-weight:750;margin:15px 0 8px;padding-bottom:5px;border-bottom:1px solid var(--line);
display:flex;align-items:center;gap:8px}}
h2.sec::before{{content:"";width:4px;height:15px;border-radius:3px;background:var(--accent)}}
.kpis{{display:grid;grid-template-columns:repeat(4,1fr);gap:10px;margin-top:13px}}
.tile{{border:1px solid var(--line);border-radius:11px;padding:9px 12px;background:var(--surface-2);
position:relative;overflow:hidden}}
.tile::before{{content:"";position:absolute;left:0;top:0;bottom:0;width:3px;background:var(--accent)}}
.tile .label{{margin-bottom:5px}}
.tile .v{{font-size:23px;font-weight:760;line-height:1}}
.tile .v .cur{{font-size:15px;font-weight:600;color:var(--muted)}}
.tile .v small{{font-size:12px;font-weight:600;color:var(--muted)}}
.tile .m{{margin-top:5px;font-size:11px;color:var(--muted)}}
.tag{{display:inline-block;font-size:10.5px;font-weight:700;padding:1px 7px;border-radius:20px;margin-left:2px}}
.tag.good{{color:var(--good);background:var(--good-soft)}}.tag.bad{{color:var(--bad);background:var(--bad-soft)}}
.exec{{margin-top:12px;padding:11px 15px;border-radius:11px;background:var(--accent-soft);
border-left:4px solid var(--accent);font-size:13px;line-height:1.55}}
.exec b{{font-weight:750}}
.cols2{{display:grid;grid-template-columns:1fr 1fr;gap:14px;margin-top:4px}}
.funnel{{display:flex;flex-direction:column;gap:2px;margin-top:8px}}
.fstep{{padding:7px 12px;border-radius:9px;color:#fff;display:flex;align-items:baseline;justify-content:space-between;
background:linear-gradient(90deg,var(--accent),var(--accent-ink))}}
.fstep .fl{{font-size:12px;font-weight:600;opacity:.95}}
.fstep .fv{{font-size:16px;font-weight:750}}
.fstep .fv small{{font-size:10px;font-weight:600;opacity:.8;margin-left:3px}}
.qrow{{display:flex;align-items:center;gap:11px;margin-top:9px}}
.qlab{{width:80px;font-size:12px;font-weight:600;flex-shrink:0}}
.qbar{{flex:1;height:8px;border-radius:6px;background:var(--line-soft);overflow:hidden}}
.qbar i{{display:block;height:100%;border-radius:6px;background:linear-gradient(90deg,var(--accent-ink),var(--accent))}}
.qpct{{width:34px;text-align:right;font-size:12px;font-weight:700;color:var(--muted)}}
table{{width:100%;border-collapse:collapse;margin-top:8px}}
thead tr{{background:var(--dark)}}
th{{text-align:left;font-size:10px;font-weight:700;letter-spacing:.05em;color:#E6EBF2;
padding:7px 10px}}
th.r,td.r{{text-align:right}}
td{{padding:8px 10px;border-bottom:1px solid var(--line-soft);font-size:13px}}
tbody tr:nth-child(even){{background:var(--surface-2)}}
tbody tr:last-child td{{border-bottom:none}}
.pname{{font-weight:650}}.pmeta{{font-size:11px;color:var(--muted);margin-top:1px}}
.big{{font-weight:700;font-size:14px}}
.pill{{display:inline-block;font-size:10.5px;font-weight:700;padding:2px 8px;border-radius:20px}}
.pill.good{{color:var(--good);background:var(--good-soft)}}.pill.warn{{color:var(--warn);background:var(--warn-soft)}}
.pill.bad{{color:var(--bad);background:var(--bad-soft)}}
.ga{{border:1px solid var(--line);border-radius:11px;padding:10px 13px}}
.ga h3{{margin:0 0 8px;font-size:12.5px;font-weight:750}}
.ga.good h3{{color:var(--good)}}.ga.bad h3{{color:var(--warn)}}
.ga ul{{margin:0;padding:0;list-style:none;display:flex;flex-direction:column;gap:8px}}
.ga li{{font-size:12.5px;line-height:1.45;color:var(--muted);padding-left:18px;position:relative}}
.ga li b{{color:var(--ink);font-weight:650}}
.ga li::before{{position:absolute;left:0;top:-1px;font-weight:800}}
.ga.good li::before{{content:"✓";color:var(--good)}}.ga.bad li::before{{content:"!";color:var(--warn)}}
.next{{margin-top:8px;padding:10px 14px;border-left:4px solid var(--accent);border-radius:0 9px 9px 0;
background:var(--accent-soft);font-size:12.5px;line-height:1.5;color:var(--ink)}}
.next b{{color:var(--accent-ink)}}
.rfoot{{margin-top:12px;padding:9px 13mm 0;border-top:1px solid var(--line);
display:flex;justify-content:space-between;gap:12px;font-size:10.5px;color:var(--faint)}}
.rfoot b{{color:var(--muted);font-weight:650}}
@media print{{@page{{size:A4;margin:0}}html,body{{background:#fff}}.bar{{display:none}}
.page{{width:auto;min-height:auto;margin:0;box-shadow:none}}
.avoid{{break-inside:avoid;page-break-inside:avoid}}}}
"""

def _bold(s):
    out, i = [], 0
    while True:
        a = s.find("**", i)
        if a < 0: out.append(s[i:]); break
        b = s.find("**", a + 2)
        if b < 0: out.append(s[i:]); break
        out.append(s[i:a]); out.append("<b>" + s[a+2:b] + "</b>"); i = b + 2
    return "".join(out)

def _wow_tag(d):
    ch = pct_change(d.get("msgs"), d.get("msgs_prev"))
    if ch is None: return ""
    return f'{"▲" if ch >= 0 else "▼"} {abs(ch):.0f}% จากช่วงก่อน'

def _cpr_tag(d):
    cpr, t = d.get("cpr"), d.get("cpr_target", 150)
    if cpr is None: return ""
    return '<span class="tag good">ผ่าน</span>' if cpr <= t else '<span class="tag bad">เกินเป้า</span>'

def build_html(d):
    period = d.get("period", "monthly")
    kpis = f"""
  <div class="kpis avoid">
    <div class="tile"><div class="label">แชท · ว่าที่ลูกค้า</div>
      <div class="v num">{d.get('msgs','—')} <small>คน</small></div><div class="m">{_wow_tag(d)}</div></div>
    <div class="tile"><div class="label">เลดจริง · คุยต่อ</div>
      <div class="v num">{d.get('leads') or '—'} <small>คน</small></div>
      <div class="m">{(str(d['lead_rate'])+'% ของแชท') if d.get('lead_rate') else '—'}
        {'<span class="tag good">คุณภาพดี</span>' if (d.get('lead_rate') or 0) >= 50 else ''}</div></div>
    <div class="tile"><div class="label">ค่าต่อ 1 แชท (CPR)</div>
      <div class="v num"><span class="cur">฿</span>{baht(d.get('cpr'))}</div>
      <div class="m">เป้า &lt; ฿{int(d.get('cpr_target',150))} {_cpr_tag(d)}</div></div>
    <div class="tile"><div class="label">ค่าแอดที่จ่าย</div>
      <div class="v num"><span class="cur">฿</span>{baht(d.get('spend'))}</div>
      <div class="m">{('จากงบ ฿'+baht(d['budget'])) if d.get('budget') else 'Reach '+baht(d.get('reach'))+' คน'}</div></div>
  </div>"""

    sec_metrics = f'<p style="font-size:11px;color:var(--muted);margin:9px 0 0">' \
        f'Reach {baht(d.get("reach"))} คน · CPM ฿{baht(d.get("cpm"))} · ความถี่ {d.get("frequency","—")}x' \
        f'{" · นัดคุย/ดูแบบ "+str(d["appointments"])+" นัด" if d.get("appointments") else ""}</p>'

    exec_html = f'<div class="exec avoid">📌 <b>สรุปผู้บริหาร:</b> {_bold(d.get("exec_summary",""))}</div>' \
        if d.get("exec_summary") else ""

    fq = ""
    if d.get("funnel") or d.get("questions"):
        fsteps = ""
        fmax = max([v for _, v in d.get("funnel", [])] or [1])
        for lbl, v in d.get("funnel", []):
            fsteps += f'<div class="fstep" style="width:{max(38,int(v/fmax*100))}%"><span class="fl">{lbl}</span>' \
                      f'<span class="fv num">{v}<small>คน</small></span></div>'
        qrows = "".join(
            f'<div class="qrow"><span class="qlab">{lbl}</span>'
            f'<span class="qbar"><i style="width:{p}%"></i></span>'
            f'<span class="qpct num">{p}%</span></div>' for lbl, p in d.get("questions", []))
        fq = f"""
  <h2 class="sec">คุณภาพแชท &amp; สิ่งที่ลูกค้าถาม</h2>
  <div class="cols2 avoid">
    <div><div class="funnel">{fsteps}</div>
      <p style="font-size:10.5px;color:var(--faint);margin:9px 0 0">วัดจากข้อความในแชท · ปิดการขายจริงทีมยืนยันเอง</p></div>
    <div>{qrows}</div></div>"""

    ptable = ""
    if d.get("personas"):
        color = {"good": "var(--good)", "warn": "var(--warn)", "bad": "var(--bad)", "no-data": "var(--muted)"}
        pill = {"good": '<span class="pill good">ผ่านเป้า</span>', "warn": '<span class="pill warn">เฝ้าระวัง</span>',
                "bad": '<span class="pill bad">เกินเป้า</span>', "no-data": '<span class="pill warn">ไม่มีข้อมูล</span>'}
        rows = ""
        for p in d["personas"]:
            st = p.get("status", "no-data")
            rows += f'<tr><td><div class="pname">{p["name"]}</div><div class="pmeta">{p.get("meta","")}</div></td>' \
                    f'<td class="r num">฿{baht(p.get("spend"))}</td><td class="r num">{p.get("msgs","—")}</td>' \
                    f'<td class="r num big" style="color:{color.get(st)}">฿{baht(p.get("cpr"))}</td>' \
                    f'<td class="r">{pill.get(st,"")}</td></tr>'
        ptable = f"""
  <h2 class="sec">ผลงานรายกลุ่มเป้าหมาย</h2>
  <table class="avoid"><thead><tr><th>กลุ่มเป้าหมาย</th><th class="r">ใช้จ่าย</th>
    <th class="r">แชท</th><th class="r">CPR</th><th class="r">สถานะ</th></tr></thead><tbody>{rows}</tbody></table>"""

    gi = ""
    if d.get("goods") or d.get("improves"):
        gl = "".join(f"<li>{_bold(x)}</li>" for x in d.get("goods", []))
        il = "".join(f"<li>{_bold(x)}</li>" for x in d.get("improves", []))
        gi = f"""
  <h2 class="sec">สิ่งที่ดี &amp; สิ่งที่ต้องปรับ</h2>
  <div class="cols2 avoid">
    <div class="ga good"><h3>✅ สิ่งที่ทำได้ดี</h3><ul>{gl}</ul></div>
    <div class="ga bad"><h3>⚠️ สิ่งที่ต้องปรับ</h3><ul>{il}</ul></div></div>"""

    nxt = f'<div class="next avoid"><b>แผน{"เดือน" if period=="monthly" else "สัปดาห์"}หน้า:</b> {_bold(d.get("next_plan",""))}</div>' \
        if d.get("next_plan") else ""

    if d.get("_mock"):
        banner = '<div class="bar"><p>⚠️ <b>ตัวอย่าง (mockup)</b> — เลขสมมติ ยังไม่ใช่ข้อมูลจริง · ' \
                 'เติมเลขจริงก่อนส่งหัวหน้า</p><span class="sp"></span>' \
                 '<button class="pbtn" onclick="window.print()">พิมพ์ / บันทึก PDF</button></div>'
    else:
        banner = '<div class="bar"><p>รายงานพร้อมส่ง · กดพิมพ์เพื่อบันทึกเป็น PDF</p>' \
                 '<span class="sp"></span><button class="pbtn" onclick="window.print()">พิมพ์ / บันทึก PDF</button></div>'

    initial = (d.get("brand", "?")[:1]) or "?"
    return f"""<!doctype html><html lang="th"><head><meta charset="utf-8">
<title>รายงานผลโฆษณา Facebook — {d.get('brand','')}</title><style>{_css()}</style></head><body>
{banner}
<div class="page">
  <div class="ribbon">
    <div class="rmark">{initial}</div>
    <div><h1>{d.get('brand','')}</h1><p class="rt">{d.get('tagline','')} · รายงานผลโฆษณา Facebook</p></div>
    <div class="rmeta"><div><b>{d.get('period_label','')}</b></div>
      <div>ช่วงข้อมูล {d.get('range_label','')}</div>
      <div><span class="chip">KPI หลัก · CPR</span></div></div>
  </div>
  <div class="body">
  {kpis}
  {sec_metrics}
  {exec_html}
  {fq}
  {ptable}
  {gi}
  {nxt}
  </div>
  <div class="rfoot">
    <span>ที่มา: Facebook Marketing API · วิเคราะห์โดย my-ads (Claude){' · <b>ตัวอย่าง</b>' if d.get('_mock') else ''}</span>
    <span>จัดทำโดย <b>ทีม {d.get('brand','')}</b> · my-ads</span>
  </div>
</div></body></html>"""

# ================================================================
# PDF (Chrome headless)
# ================================================================
def _find_chrome():
    """หา Chrome/Edge/Chromium ข้าม OS (Mac / Windows / Linux)"""
    from shutil import which
    la = os.environ.get("LOCALAPPDATA", "")
    pf = os.environ.get("PROGRAMFILES", r"C:\Program Files")
    pf86 = os.environ.get("PROGRAMFILES(X86)", r"C:\Program Files (x86)")
    candidates = [
        # macOS
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        # Windows (Chrome + Edge fallback — Edge เป็น Chromium ใช้ --print-to-pdf ได้)
        os.path.join(pf, r"Google\Chrome\Application\chrome.exe"),
        os.path.join(pf86, r"Google\Chrome\Application\chrome.exe"),
        os.path.join(la, r"Google\Chrome\Application\chrome.exe") if la else "",
        os.path.join(pf86, r"Microsoft\Edge\Application\msedge.exe"),
        os.path.join(pf, r"Microsoft\Edge\Application\msedge.exe"),
        # Linux / PATH
        "google-chrome", "google-chrome-stable", "chromium", "chromium-browser", "microsoft-edge",
    ]
    for c in candidates:
        if not c:
            continue
        if os.path.exists(c) or which(c):
            return c
    return None

def build_pdf(html_path, pdf_path):
    chrome = _find_chrome()
    if not chrome:
        print("⚠️ ไม่เจอ Chrome/Edge — เปิดไฟล์ HTML ในเบราว์เซอร์แล้วกด Ctrl/Cmd+P → บันทึกเป็น PDF แทนได้")
        return None
    try:
        subprocess.run([chrome, "--headless", "--disable-gpu", "--no-pdf-header-footer",
                        f"--print-to-pdf={pdf_path}", _file_uri(html_path)],
                       check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except Exception as e:
        print(f"⚠️ สร้าง PDF ไม่สำเร็จ ({e}) — เปิด HTML แล้วพิมพ์เป็น PDF เองได้"); return None
    return pdf_path

def _file_uri(path):
    ap = os.path.abspath(path)
    return "file:///" + ap.replace("\\", "/") if os.name == "nt" else "file://" + ap

# ================================================================
# PPTX (มาตรฐานเอเจนซี่: ribbon + accent CI + cover + exec)
# ================================================================
def build_pptx(d, pptx_path):
    try:
        from pptx import Presentation
    except ImportError:
        exe = os.path.basename(sys.executable) or "python3"
        print("⚠️ ยังไม่มี python-pptx (ตัวสร้าง PPTX) — ติดตั้งด้วย:")
        print(f"     {exe} -m pip install python-pptx")
        print("   (Windows ถ้า python ใช้ไม่ได้ ลอง:  py -m pip install python-pptx)")
        print("   ข้าม PPTX ไปก่อน · HTML/PDF ยังออกได้ปกติ")
        return None
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    def rgb(h): return RGBColor.from_string(h)
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    W, H = 13.333, 7.5
    blank = prs.slide_layouts[6]
    period = d.get("period", "monthly")

    def slide():
        s = prs.slides.add_slide(blank)
        s.background.fill.solid(); s.background.fill.fore_color.rgb = rgb(PAPER)
        return s

    def txt(s, x, y, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT,
            anchor=MSO_ANCHOR.TOP, font=TH_FONT):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
        p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = 1.1
        for chunk, b in _split_bold(text):
            r = p.add_run(); r.text = chunk
            r.font.name = font; r.font.size = Pt(size)
            r.font.bold = bold or b; r.font.color.rgb = rgb(color)
        return tb

    def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=False):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
        shp.shadow.inherit = False
        if fill: shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
        else: shp.fill.background()
        if line: shp.line.color.rgb = rgb(line); shp.line.width = Pt(line_w)
        else: shp.line.fill.background()
        return shp

    def ribbon(s, title, subtitle):
        rect(s, 0, 0, W, 1.0, fill=DARK)
        rect(s, 0, 1.0, W, 0.045, fill=ACCENT)
        rect(s, 0.5, 0.24, 0.55, 0.55, fill=ACCENT, radius=True)
        txt(s, 0.5, 0.24, 0.55, 0.55, (d.get("brand","?")[:1]), size=22, bold=True,
            color="FFFFFF", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, 1.25, 0.16, 8.5, 0.5, title, size=21, bold=True, color="FFFFFF")
        txt(s, 1.25, 0.62, 8.5, 0.32, subtitle, size=11.5, color="C7D0DD")
        txt(s, W - 4.3, 0.2, 3.8, 0.3, d.get("period_label",""), size=11.5, bold=True,
            color="E6EBF2", align=PP_ALIGN.RIGHT)
        txt(s, W - 4.3, 0.56, 3.8, 0.3, "ช่วงข้อมูล " + d.get("range_label",""), size=10,
            color="AAB4C4", align=PP_ALIGN.RIGHT)

    def foot(s):
        txt(s, 0.5, H - 0.42, W - 1, 0.3,
            ("ตัวอย่าง (mockup) · " if d.get("_mock") else "") +
            f"my-ads · Facebook Marketing API · {d.get('period_label','')}",
            size=9, color=FAINT, align=PP_ALIGN.CENTER)

    def kpi_card(s, x, y, w, h, label, value, sub, vcolor=ACCENT):
        rect(s, x, y, w, h, fill=SURFACE, line=LINE, line_w=1.0, radius=True)
        rect(s, x, y + 0.12, 0.06, h - 0.24, fill=vcolor)  # accent bar ซ้าย
        txt(s, x + 0.22, y + 0.13, w - 0.35, 0.3, label, size=10, bold=True, color=FAINT)
        txt(s, x + 0.2, y + 0.42, w - 0.35, h - 0.7, value, size=27, bold=True, color=vcolor,
            font=EN_FONT)
        if sub:
            txt(s, x + 0.22, y + h - 0.35, w - 0.4, 0.3, sub, size=9.5, color=MUTED)

    # ---- Slide 1: COVER ----
    s = slide()
    rect(s, 0, 0, 0.28, H, fill=DARK)
    rect(s, 0.28, 0, 0.05, H, fill=ACCENT)
    rect(s, 0.9, 1.15, 0.95, 0.95, fill=ACCENT, radius=True)
    txt(s, 0.9, 1.15, 0.95, 0.95, (d.get("brand","?")[:1]), size=40, bold=True, color="FFFFFF",
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 2.05, 1.2, 9, 0.7, d.get("brand",""), size=32, bold=True, color=INK)
    txt(s, 2.05, 1.98, 9, 0.4, d.get("tagline",""), size=15, color=MUTED)
    txt(s, 0.9, 2.75, 11, 0.6, "รายงานผลโฆษณา Facebook", size=22, bold=True, color=ACCENT_INK)
    txt(s, 0.9, 3.42, 11, 0.4, d.get("period_label","") + "   ·   " + d.get("range_label",""),
        size=14, bold=True, color=INK)
    rect(s, 0.92, 4.0, 4.5, 0.04, fill=ACCENT)
    # quick stat cards
    stats = [
        ("แชท · ว่าที่ลูกค้า", f"{d.get('msgs','—')}", "คน" + ("  " + _wow_tag(d) if _wow_tag(d) else "")),
        ("ค่าต่อ 1 แชท (CPR)", f"฿{baht(d.get('cpr'))}", f"เป้า < ฿{int(d.get('cpr_target',150))}"),
        ("ค่าแอดที่จ่าย", f"฿{baht(d.get('spend'))}",
         (f"จากงบ ฿{baht(d['budget'])}" if d.get('budget') else f"Reach {baht(d.get('reach'))}")),
    ]
    cw = 3.7; y0 = 4.55
    for i, (lb, v, sub) in enumerate(stats):
        vc = ACCENT
        if i == 1 and d.get("cpr") is not None:
            vc = GOOD if d["cpr"] <= d.get("cpr_target", 150) else BAD
        kpi_card(s, 0.9 + i * (cw + 0.28), y0, cw, 1.55, lb, v, sub, vcolor=vc)
    txt(s, 0.9, H - 0.55, 11, 0.35, ("ตัวอย่าง (mockup) — เลขสมมติ · " if d.get("_mock") else "") +
        "จัดทำโดยทีม " + d.get("brand","") + " · my-ads", size=10, color=FAINT)

    # ---- Slide 2: EXECUTIVE SUMMARY (band + text) ----
    s = slide()
    ribbon(s, "สรุปผู้บริหาร", "Executive Summary · ภาพรวมเดือน")
    band = [
        ("แชท · ว่าที่ลูกค้า", f"{d.get('msgs','—')}", "คน" + (" " + _wow_tag(d) if _wow_tag(d) else ""), ACCENT),
        ("เลดจริง · คุยต่อ", f"{d.get('leads') or '—'}",
         (f"{d['lead_rate']}% ของแชท" if d.get('lead_rate') else "—"),
         GOOD if (d.get('lead_rate') or 0) >= 50 else ACCENT),
        ("ค่าต่อ 1 แชท (CPR)", f"฿{baht(d.get('cpr'))}", f"เป้า < ฿{int(d.get('cpr_target',150))}",
         GOOD if (d.get('cpr') or 9e9) <= d.get('cpr_target', 150) else BAD),
        ("ค่าแอดที่จ่าย", f"฿{baht(d.get('spend'))}",
         (f"จากงบ ฿{baht(d['budget'])}" if d.get('budget') else f"Reach {baht(d.get('reach'))}"), ACCENT),
    ]
    cw = 2.95
    for i, (lb, v, sub, vc) in enumerate(band):
        kpi_card(s, 0.5 + i * (cw + 0.1), 1.35, cw, 1.35, lb, v, sub, vcolor=vc)
    txt(s, 0.5, 2.92, 12.3, 0.3,
        f"Reach {baht(d.get('reach'))} คน · CPM ฿{baht(d.get('cpm'))} · ความถี่ {d.get('frequency','—')}x"
        + (f" · นัดคุย {d['appointments']} นัด" if d.get('appointments') else ""),
        size=10.5, color=MUTED)
    if d.get("exec_summary"):
        rect(s, 0.5, 3.4, 12.33, 2.5, fill=ACCENT_SOFT, radius=True)
        rect(s, 0.5, 3.4, 0.08, 2.5, fill=ACCENT)
        txt(s, 0.85, 3.6, 11.7, 0.4, "📌 สรุปสำหรับผู้บริหาร", size=14, bold=True, color=ACCENT_INK)
        txt(s, 0.85, 4.1, 11.7, 1.7, d.get("exec_summary",""), size=14, color=INK)
    foot(s)

    # ---- Slide 3: funnel + questions (monthly, if data) ----
    if period == "monthly" and (d.get("funnel") or d.get("questions")):
        s = slide()
        ribbon(s, "คุณภาพแชท & สิ่งที่ลูกค้าถาม", "Chat Quality & Top Questions")
        fmax = max([v for _, v in d.get("funnel", [])] or [1]); fy = 1.55
        txt(s, 0.5, 1.15, 6, 0.35, "เส้นทางแชท → นัด", size=13, bold=True, color=ACCENT_INK)
        for lbl, v in d.get("funnel", []):
            w = max(3.0, v / fmax * 5.7)
            rect(s, 0.5, fy, w, 0.62, fill=ACCENT, radius=True)
            txt(s, 0.7, fy, w - 0.3, 0.62, lbl, size=12, bold=True, color="FFFFFF", anchor=MSO_ANCHOR.MIDDLE)
            txt(s, 0.5, fy, 5.7, 0.62, f"{v} คน  ", size=13, bold=True, color="FFFFFF",
                align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, font=EN_FONT)
            fy += 0.82
        qx, qy = 7.2, 1.55
        txt(s, qx, 1.15, 5, 0.35, "ลูกค้าถามอะไรมากสุด", size=13, bold=True, color=ACCENT_INK)
        for lbl, p in d.get("questions", []):
            txt(s, qx, qy, 1.9, 0.35, lbl, size=12)
            rect(s, qx + 2.0, qy + 0.07, 2.7, 0.2, fill=LINE_SOFT, radius=True)
            rect(s, qx + 2.0, qy + 0.07, max(0.1, 2.7 * p / 100), 0.2, fill=ACCENT, radius=True)
            txt(s, qx + 4.8, qy, 0.7, 0.35, f"{p}%", size=12, bold=True, color=MUTED, font=EN_FONT)
            qy += 0.52
        foot(s)

    # ---- Slide 4: persona table ----
    if d.get("personas"):
        s = slide()
        ribbon(s, "ผลงานรายกลุ่มเป้าหมาย", "Performance by Audience")
        rows = d["personas"][:8]
        from pptx.util import Inches as In
        tbl = s.shapes.add_table(len(rows) + 1, 5, In(0.5), In(1.35), In(12.33),
                                 In(0.52 * (len(rows) + 1))).table
        for j, w in enumerate([5.0, 2.1, 1.5, 1.8, 1.93]): tbl.columns[j].width = In(w)
        heads = ["กลุ่มเป้าหมาย", "ใช้จ่าย", "แชท", "CPR", "สถานะ"]
        stat_txt = {"good": "ผ่านเป้า", "warn": "เฝ้าระวัง", "bad": "เกินเป้า", "no-data": "ไม่มีข้อมูล"}
        stat_col = {"good": GOOD, "warn": WARN, "bad": BAD, "no-data": MUTED}
        for j, h in enumerate(heads):
            _cell(tbl.cell(0, j), h, 11, True, "FFFFFF", DARK,
                  PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT, EN_FONT if j else TH_FONT)
        for i, p in enumerate(rows, 1):
            st = p.get("status", "no-data")
            vals = [p["name"], f"฿{baht(p.get('spend'))}", str(p.get("msgs", "—")),
                    f"฿{baht(p.get('cpr'))}", stat_txt.get(st, "")]
            for j, v in enumerate(vals):
                col = stat_col.get(st) if j in (3, 4) else INK
                _cell(tbl.cell(i, j), v, 12, j in (3, 4), col, SURFACE if i % 2 else SURFACE2,
                      PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT, EN_FONT if j in (1, 2, 3) else TH_FONT)
        foot(s)

    # ---- Slide 5: good / improve + next ----
    if d.get("goods") or d.get("improves") or d.get("next_plan"):
        s = slide()
        ribbon(s, "สรุป & แผนต่อไป", "Wins · Fixes · Next Plan")
        rect(s, 0.5, 1.35, 5.95, 3.5, fill=GOOD_SOFT, radius=True)
        rect(s, 0.5, 1.35, 0.07, 3.5, fill=GOOD)
        txt(s, 0.75, 1.5, 5.6, 0.4, "✅ สิ่งที่ทำได้ดี", size=14, bold=True, color=GOOD)
        yy = 2.02
        for g in d.get("goods", []):
            txt(s, 0.78, yy, 5.5, 0.75, "•  " + g, size=12, color=INK); yy += 0.72
        rect(s, 6.85, 1.35, 5.95, 3.5, fill=WARN_SOFT, radius=True)
        rect(s, 6.85, 1.35, 0.07, 3.5, fill=WARN)
        txt(s, 7.1, 1.5, 5.6, 0.4, "⚠️ สิ่งที่ต้องปรับ", size=14, bold=True, color=WARN)
        yy = 2.02
        for im in d.get("improves", []):
            txt(s, 7.13, yy, 5.5, 0.75, "•  " + im, size=12, color=INK); yy += 0.72
        if d.get("next_plan"):
            rect(s, 0.5, 5.1, 12.33, 1.55, fill=ACCENT_SOFT, radius=True)
            rect(s, 0.5, 5.1, 0.08, 1.55, fill=ACCENT)
            lbl = "แผนเดือนหน้า" if period == "monthly" else "แผนสัปดาห์หน้า"
            txt(s, 0.8, 5.25, 11.8, 0.4, "🎯 " + lbl, size=13, bold=True, color=ACCENT_INK)
            txt(s, 0.8, 5.68, 11.8, 0.9, d.get("next_plan",""), size=13, color=INK)
        foot(s)

    prs.save(pptx_path)
    return pptx_path

def _cell(cell, text, size, bold, color, fill, align, font=TH_FONT):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(7); cell.margin_top = cell.margin_bottom = Pt(3)
    tf = cell.text_frame; p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = _plain(text)
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = RGBColor.from_string(color)

def _split_bold(s):
    parts, i = [], 0
    while True:
        a = s.find("**", i)
        if a < 0: parts.append((s[i:], False)); break
        b = s.find("**", a + 2)
        if b < 0: parts.append((s[i:], False)); break
        if a > i: parts.append((s[i:a], False))
        parts.append((s[a+2:b], True)); i = b + 2
    return [p for p in parts if p[0]] or [("", False)]

def _plain(s):
    return s.replace("**", "")

# ================================================================
# main
# ================================================================
def default_range(period):
    from datetime import date, timedelta
    t = date.today()
    if period == "weekly":
        end = t - timedelta(days=1); return (end - timedelta(days=6)).isoformat(), end.isoformat()
    first = t.replace(day=1); last_prev = first - timedelta(days=1)
    return last_prev.replace(day=1).isoformat(), last_prev.isoformat()

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--period", choices=["weekly", "monthly"], default="monthly")
    ap.add_argument("--since"); ap.add_argument("--until")
    ap.add_argument("--mock", action="store_true", help="ใช้เลขตัวอย่าง ไม่ต่อ API")
    ap.add_argument("--fetch", action="store_true", help="ดึงเลขจาก API → report_data.json")
    ap.add_argument("--render", action="store_true", help="อ่าน report_data.json → ออกไฟล์")
    ap.add_argument("--out", default="all", help="html | pdf | pptx | all")
    a = ap.parse_args()

    since, until = a.since, a.until
    if not (since and until):
        since, until = default_range(a.period)

    if a.mock:
        data = mock(a.period)
    elif a.fetch:
        data = fetch(a.period, since, until)
        json.dump(data, open(DATA_JSON, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
        print(f"✅ ดึงตัวเลขแล้ว → {DATA_JSON}")
        print("👉 ต่อไป: Claude เติมช่อง exec_summary / funnel / questions / goods / improves / next_plan")
        print("   แล้วรัน:  python3 gen_report.py --period", a.period, "--render --out all")
        return
    elif a.render:
        if not os.path.exists(DATA_JSON):
            sys.exit("❌ ไม่เจอ report_data.json — รัน --fetch ก่อน (หรือ --mock เพื่อดูตัวอย่าง)")
        data = json.load(open(DATA_JSON, encoding="utf-8"))
    else:
        sys.exit("เลือกโหมด: --mock (ดูตัวอย่าง) | --fetch (ดึงเลข) | --render (ออกไฟล์)")

    prefix = f"รายงานแอด-{'เดือน' if a.period=='monthly' else 'สัปดาห์'}"
    outs = ["html", "pdf", "pptx"] if a.out == "all" else [a.out]
    made = []
    html_path = os.path.join(HERE, prefix + ".html")
    if "html" in outs or "pdf" in outs:
        open(html_path, "w", encoding="utf-8").write(build_html(data))
        if "html" in outs: made.append(html_path)
    if "pdf" in outs:
        pdf = build_pdf(html_path, os.path.join(HERE, prefix + ".pdf"))
        if pdf: made.append(pdf)
    if "pptx" in outs:
        made.append(build_pptx(data, os.path.join(HERE, prefix + ".pptx")))

    print("✅ ออกรายงานแล้ว:")
    for m in made:
        print("   •", m)
    if data.get("_mock"):
        print("⚠️  นี่คือตัวอย่าง (เลขสมมติ) — ใช้ --fetch เพื่อดึงเลขจริงก่อนส่งหัวหน้า")

if __name__ == "__main__":
    main()
